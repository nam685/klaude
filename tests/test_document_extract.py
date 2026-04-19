"""Tests for the shared document extraction helpers."""

import base64
import shutil
import subprocess
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest

from tests.fixtures import make_docx, make_pptx, make_xlsx

from klaude.config import VisionConfig
from klaude.tools._document import (
    MAX_EXTRACTED_BYTES,
    extract,
    _apply_cap,
    _wrap,
)


def test_wrap_includes_system_reminder_and_document_tags() -> None:
    out = _wrap("hello", path="/x.txt", fmt="txt")
    assert "<system-reminder>" in out
    assert "</system-reminder>" in out
    assert 'path="/x.txt"' in out
    assert 'format="txt"' in out
    assert "hello" in out
    assert "untrusted data" in out


def test_cap_truncates_long_content() -> None:
    huge = "a" * (MAX_EXTRACTED_BYTES + 5000)
    capped = _apply_cap(huge)
    assert "[truncated at 200 KB — original document was larger]" in capped
    head = capped.split("\n\n[truncated")[0]
    assert len(head.encode("utf-8")) <= MAX_EXTRACTED_BYTES


def test_cap_leaves_short_content_alone() -> None:
    assert _apply_cap("short") == "short"


def test_extract_missing_file(tmp_path: Path) -> None:
    out = extract(tmp_path / "nope.pdf")
    assert out.startswith("Error:")
    assert "not found" in out.lower() or "no such" in out.lower()


def test_extract_unsupported_extension(tmp_path: Path) -> None:
    p = tmp_path / "foo.bogus"
    p.write_text("whatever")
    out = extract(p)
    assert out.startswith("Error:")
    assert "unsupported" in out.lower() or "not supported" in out.lower()


def test_html_strips_tags(tmp_path: Path) -> None:
    p = tmp_path / "page.html"
    p.write_text(
        "<html><body><h1>Hello</h1><p>World <b>here</b></p>"
        "<script>alert('x')</script></body></html>"
    )
    out = extract(p)
    assert "Hello" in out
    assert "World" in out
    assert "here" in out
    assert "<h1>" not in out
    assert "<script>" not in out
    # Script content dropped entirely
    assert "alert" not in out
    # Structural breaks preserved between block elements
    assert "Hello\n" in out or "Hello\n\n" in out


def test_html_preserves_paragraph_breaks(tmp_path: Path) -> None:
    p = tmp_path / "multi.html"
    p.write_text(
        "<html><body>"
        "<p>First paragraph.</p>"
        "<p>Second paragraph.</p>"
        "<h2>Heading</h2>"
        "<p>Third.</p>"
        "</body></html>"
    )
    out = extract(p)
    # Paragraphs on separate lines
    lines = [ln for ln in out.splitlines() if "First paragraph" in ln or "Second paragraph" in ln]
    assert len(lines) == 2
    # Heading present as its own line
    heading_idx = next(i for i, ln in enumerate(out.splitlines()) if "Heading" in ln)
    third_idx = next(i for i, ln in enumerate(out.splitlines()) if "Third" in ln)
    assert third_idx > heading_idx


def test_html_htm_extension(tmp_path: Path) -> None:
    p = tmp_path / "page.htm"
    p.write_text("<p>htm works</p>")
    out = extract(p)
    assert "htm works" in out


def test_docx_extracts_paragraphs(tmp_path: Path) -> None:
    p = make_docx(tmp_path / "tiny.docx", ["First line", "Second line", "Third"])
    out = extract(p)
    assert "First line" in out
    assert "Second line" in out
    assert "Third" in out
    assert 'format="docx"' in out


def test_docx_extracts_table_cells(tmp_path: Path) -> None:
    from docx import Document

    p = tmp_path / "with_table.docx"
    doc = Document()
    doc.add_paragraph("Intro paragraph")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "A1"
    table.rows[0].cells[1].text = "B1"
    table.rows[1].cells[0].text = "A2"
    table.rows[1].cells[1].text = "B2"
    doc.add_paragraph("Trailing paragraph")
    doc.save(str(p))

    out = extract(p)
    assert "Intro paragraph" in out
    assert "A1" in out and "B1" in out
    assert "A2" in out and "B2" in out
    assert "Trailing paragraph" in out


def test_xlsx_extracts_all_sheets_as_csv(tmp_path: Path) -> None:
    p = make_xlsx(
        tmp_path / "tiny.xlsx",
        {
            "People": [["name", "age"], ["Ada", 36], ["Grace", 45]],
            "Notes": [["topic"], ["compiler"]],
        },
    )
    out = extract(p)
    assert "# Sheet: People" in out
    assert "# Sheet: Notes" in out
    assert "name,age" in out
    assert "Ada,36" in out
    assert "Grace,45" in out
    assert "topic" in out
    assert "compiler" in out
    assert "\n---\n" in out


def test_xlsx_formats_dates_as_iso(tmp_path: Path) -> None:
    import datetime as dt

    p = make_xlsx(
        tmp_path / "dates.xlsx",
        {
            "Log": [
                ["when", "note"],
                [dt.datetime(2026, 4, 19, 14, 30, 0), "first"],
                [dt.date(2026, 4, 20), "second"],
            ],
        },
    )
    out = extract(p)
    assert "2026-04-19T14:30:00" in out
    assert "2026-04-20" in out
    # Should NOT contain the Python repr spelling.
    assert "2026-04-19 14:30:00" not in out


def test_pptx_extracts_slide_text(tmp_path: Path) -> None:
    p = make_pptx(
        tmp_path / "tiny.pptx",
        [
            ["Intro", "hello world"],
            ["Findings", "result: 42"],
        ],
    )
    out = extract(p)
    assert "# Slide 1" in out
    assert "# Slide 2" in out
    assert "Intro" in out
    assert "hello world" in out
    assert "Findings" in out
    assert "result: 42" in out
    assert "\n---\n" in out


def test_pptx_extracts_table_cells(tmp_path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    p = tmp_path / "table.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Results"
    shape = slide.shapes.add_table(
        rows=2, cols=2,
        left=Inches(1), top=Inches(2),
        width=Inches(6), height=Inches(2),
    )
    t = shape.table
    t.cell(0, 0).text = "metric"
    t.cell(0, 1).text = "value"
    t.cell(1, 0).text = "latency"
    t.cell(1, 1).text = "42ms"
    prs.save(str(p))
    out = extract(p)
    assert "Results" in out
    assert "metric" in out and "value" in out
    assert "latency" in out and "42ms" in out


def test_pptx_extracts_grouped_shapes(tmp_path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    p = tmp_path / "grouped.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Header"
    tb1 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(1))
    tb1.text_frame.text = "child one"
    tb2 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(3), Inches(1))
    tb2.text_frame.text = "child two"
    prs.save(str(p))
    out = extract(p)
    assert "child one" in out
    assert "child two" in out


# ---------------------------------------------------------------------------
# PDF tests
# ---------------------------------------------------------------------------

PDFTOTEXT = shutil.which("pdftotext")


@pytest.mark.skipif(PDFTOTEXT is None, reason="pdftotext not installed")
def test_pdf_extracts_text(tmp_path: Path) -> None:
    # Build a tiny PDF via pandoc if available, else skip.
    if shutil.which("pandoc") is None:
        pytest.skip("pandoc not installed")
    md = tmp_path / "in.md"
    md.write_text("hello PDF world\n")
    pdf = tmp_path / "tiny.pdf"
    subprocess.run(["pandoc", str(md), "-o", str(pdf)], check=True)
    out = extract(pdf)
    assert "hello PDF world" in out
    assert 'format="pdf"' in out


def test_pdf_missing_binary_gives_install_hint(tmp_path: Path) -> None:
    p = tmp_path / "x.pdf"
    p.write_bytes(b"%PDF-1.4\n%dummy\n")
    with patch("klaude.tools._document.shutil.which", return_value=None):
        out = extract(p)
    assert "pdftotext not found" in out
    assert "brew install poppler" in out or "apt install poppler-utils" in out


def test_pdf_encrypted_clean_error(tmp_path: Path) -> None:
    p = tmp_path / "enc.pdf"
    p.write_bytes(b"%PDF-1.4\n")

    def fake_run(*args, **_kwargs):
        return subprocess.CompletedProcess(
            args=args, returncode=3, stdout=b"", stderr=b"Error: PDF file is encrypted"
        )

    with patch("klaude.tools._document.shutil.which", return_value="/usr/bin/pdftotext"), \
         patch("klaude.tools._document.subprocess.run", side_effect=fake_run):
        out = extract(p)
    assert "password-protected" in out.lower() or "encrypted" in out.lower()


def test_ocr_missing_binary_gives_install_hint(tmp_path: Path) -> None:
    p = tmp_path / "x.png"
    p.write_bytes(b"\x89PNG\r\n\x1a\n")  # binary presence; not a real PNG
    from klaude.tools._document import _extract_image_ocr

    msg = ""
    with patch("klaude.tools._document.shutil.which", return_value=None):
        try:
            _extract_image_ocr(p)
        except RuntimeError as e:
            msg = str(e)
    assert "tesseract not found" in msg
    assert "brew install tesseract" in msg or "apt install tesseract-ocr" in msg


# ---------------------------------------------------------------------------
# Image VLM / dispatcher tests
# ---------------------------------------------------------------------------

IMAGE_EXTS = (".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif", ".webp")


def _tiny_png(path: Path) -> Path:
    # 1x1 red PNG (67 bytes) — smallest legal PNG payload.
    data = bytes.fromhex(
        "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
        "890000000d49444154789c62f8cf000000000300015f5d9b8a0000000049454e"
        "44ae426082"
    )
    path.write_bytes(data)
    return path


def test_image_vlm_backend_calls_model(tmp_path: Path, monkeypatch) -> None:
    monkeypatch.setenv("OPENROUTER_API_KEY", "sk-test")
    from klaude.tools import _document as d

    fake_client = MagicMock()
    fake_client.chat.completions.create.return_value = MagicMock(
        choices=[MagicMock(message=MagicMock(content="a red pixel"))]
    )
    monkeypatch.setattr(d, "_openai_client", lambda _cfg: fake_client)

    p = _tiny_png(tmp_path / "tiny.png")
    cfg = VisionConfig()
    monkeypatch.setattr(d, "_vision_config", lambda: cfg)

    out = d.extract(p)
    assert "a red pixel" in out
    assert 'format="png"' in out

    call = fake_client.chat.completions.create.call_args
    messages = call.kwargs["messages"]
    parts = messages[0]["content"]
    image_part = next(part for part in parts if part["type"] == "image_url")
    url = image_part["image_url"]["url"]
    assert url.startswith("data:image/png;base64,")
    b64 = url.split(",", 1)[1]
    assert base64.b64decode(b64) == p.read_bytes()


def test_image_vlm_fallback_noted_when_key_unset(tmp_path: Path, monkeypatch) -> None:
    monkeypatch.delenv("OPENROUTER_API_KEY", raising=False)
    from klaude.tools import _document as d

    monkeypatch.setattr(d, "_extract_image_ocr", lambda _p: "ocr text here")
    monkeypatch.setattr(
        d,
        "_vision_config",
        lambda: VisionConfig(backend="vlm", fallback="ocr",
                             api_key_env="OPENROUTER_API_KEY"),
    )

    p = _tiny_png(tmp_path / "tiny.png")
    out = d.extract(p)
    assert "[vision.backend=vlm but $OPENROUTER_API_KEY unset; used OCR fallback]" in out
    assert "ocr text here" in out


def test_image_vlm_fallback_error_when_configured(tmp_path: Path, monkeypatch) -> None:
    monkeypatch.delenv("OPENROUTER_API_KEY", raising=False)
    from klaude.tools import _document as d
    monkeypatch.setattr(
        d,
        "_vision_config",
        lambda: VisionConfig(backend="vlm", fallback="error",
                             api_key_env="OPENROUTER_API_KEY"),
    )
    monkeypatch.setattr(d, "_extract_image_ocr", lambda _p: pytest.fail("OCR called"))

    p = _tiny_png(tmp_path / "tiny.png")
    out = d.extract(p)
    assert out.startswith("Error:")
    assert "OPENROUTER_API_KEY" in out


def test_image_backend_ocr_direct(tmp_path: Path, monkeypatch) -> None:
    from klaude.tools import _document as d
    monkeypatch.setattr(d, "_extract_image_ocr", lambda _p: "plain ocr")
    monkeypatch.setattr(
        d, "_vision_config", lambda: VisionConfig(backend="ocr"),
    )
    p = _tiny_png(tmp_path / "tiny.png")
    out = d.extract(p)
    assert "plain ocr" in out
    assert "used OCR fallback" not in out


def test_image_extensions_all_dispatched(tmp_path: Path, monkeypatch) -> None:
    from klaude.tools import _document as d
    monkeypatch.setattr(d, "_extract_image_ocr", lambda p: f"ocr:{p.suffix}")
    monkeypatch.setattr(
        d, "_vision_config", lambda: VisionConfig(backend="ocr"),
    )
    for ext in IMAGE_EXTS:
        p = _tiny_png(tmp_path / f"f{ext}")
        out = d.extract(p)
        assert f"ocr:{ext}" in out, f"{ext} not dispatched"


def test_image_vlm_key_inherited_from_llm(tmp_path: Path, monkeypatch) -> None:
    # Primary LLM config declares api_key_env=OPENROUTER_API_KEY; vision has no
    # explicit api_key_env. Ensure the VLM path picks up the inherited env var.
    (tmp_path / ".klaude.toml").write_text(
        '[default]\nmodel = "remote"\napi_key_env = "OPENROUTER_API_KEY"\n'
    )
    monkeypatch.setenv("OPENROUTER_API_KEY", "sk-inherited")
    monkeypatch.chdir(tmp_path)

    from klaude.tools import _document as d

    seen_api_key: list[str] = []

    def fake_client(cfg):
        seen_api_key.append(cfg.api_key_env)
        client = MagicMock()
        client.chat.completions.create.return_value = MagicMock(
            choices=[MagicMock(message=MagicMock(content="inherited key ok"))]
        )
        return client

    monkeypatch.setattr(d, "_openai_client", fake_client)

    p = _tiny_png(tmp_path / "tiny.png")
    out = d.extract(p)
    assert "inherited key ok" in out
    assert seen_api_key == ["OPENROUTER_API_KEY"]
