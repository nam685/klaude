"""Fixture builders — generate tiny Office documents at test time."""

from pathlib import Path


def make_docx(path: Path, paragraphs: list[str]) -> Path:
    from docx import Document

    doc = Document()
    for p in paragraphs:
        doc.add_paragraph(p)
    doc.save(str(path))
    return path


def make_xlsx(path: Path, sheets: dict[str, list[list[object]]]) -> Path:
    from openpyxl import Workbook

    wb = Workbook()
    wb.remove(wb.active)
    for name, rows in sheets.items():
        ws = wb.create_sheet(title=name)
        for row in rows:
            ws.append(row)
    wb.save(str(path))
    return path


def make_pptx(path: Path, slides: list[list[str]]) -> Path:
    from pptx import Presentation

    prs = Presentation()
    blank = prs.slide_layouts[5]  # title only
    for texts in slides:
        slide = prs.slides.add_slide(blank)
        slide.shapes.title.text = texts[0] if texts else ""
        for extra in texts[1:]:
            left = top = 100000
            width = height = 5000000
            tb = slide.shapes.add_textbox(left, top, width, height)
            tb.text_frame.text = extra
    prs.save(str(path))
    return path
